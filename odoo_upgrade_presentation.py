import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from datetime import datetime

# Create a new presentation
presentation = Presentation()

# Function to add images to slides
def add_image(slide, image_path, left, top, width, height):
    pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    return pic

# Add a rectangle shape with text
def add_rectangle_with_text(slide, left, top, width, height, text, fill_color=None, text_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if text_color:
        shape.text_frame.text = text
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = text_color
    return shape

# Title slide
title_slide_layout = presentation.slide_layouts[0]
slide = presentation.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Upgrade Process for Odoo Database"
subtitle.text = "Using the Python-pptx Library"

# Slide 1: Introduction
slide1_layout = presentation.slide_layouts[1]
slide1 = presentation.slides.add_slide(slide1_layout)
slide1.shapes.title.text = "Introduction"
slide1.placeholders[1].text = (
    "- What is Odoo?\n"
    "- Why upgrade the Odoo database?\n"
    "- What is Python-pptx?"
)

# Slide 2: Prerequisites
slide2_layout = presentation.slide_layouts[1]
slide2 = presentation.slides.add_slide(slide2_layout)
slide2.shapes.title.text = "Prerequisites"
slide2.placeholders[1].text = (
    "- Python and Odoo installation\n"
    "- Python-pptx library\n"
    "- Backup your Odoo database\n"
    "- Latest Odoo version"
)

# Slide 3: Python-pptx Overview
slide3_layout = presentation.slide_layouts[1]
slide3 = presentation.slides.add_slide(slide3_layout)
slide3.shapes.title.text = "Python-pptx Overview"
add_rectangle_with_text(slide3, 
                        Inches(3), 
                        Inches(2),
                        Inches(4),
                        Inches(0.7), 
                        "- Installing the library\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide3, 
                        Inches(3), 
                        Inches(3), 
                        Inches(4), 
                        Inches(0.7), 
                        "- Basic syntax and usage\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide3, 
                        Inches(3), 
                        Inches(4), 
                        Inches(4), 
                        Inches(0.8), 
                        "- Creating a PowerPoint presentation with Python",
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))

# Slide 4: Odoo Database Upgrade Process Overview
slide4_layout = presentation.slide_layouts[1]
slide4 = presentation.slides.add_slide(slide4_layout)
slide4.shapes.title.text = "Odoo Database Upgrade Process Overview"
slide4.placeholders[1].text = (
    "- Exporting the Odoo database\n"
    "- Uploading the database to the Odoo Upgrade Platform\n"
    "- Testing the upgraded database\n"
    "- Migrating the database to the production environment"
)

# Slide 5: Exporting the Odoo Database (with shapes)
slide5_layout = presentation.slide_layouts[1]
slide5 = presentation.slides.add_slide(slide5_layout)
slide5.shapes.title.text = "Exporting the Odoo Database"

add_rectangle_with_text(slide5, 
                        Inches(1), 
                        Inches(2),
                        Inches(4),
                        Inches(0.5), 
                        "1. Navigate to the Odoo Database Manager", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide5, 
                        Inches(1), 
                        Inches(2.7), 
                        Inches(4), 
                        Inches(0.5), 
                        "2. Select the database to upgrade", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide5, 
                        Inches(1), 
                        Inches(3.4), 
                        Inches(4), 
                        Inches(0.5), 
                        "3. Click on the 'Backup' button", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide5, 
                        Inches(1), 
                        Inches(4.1), 
                        Inches(4), 
                        Inches(0.5), 
                        "4. Choose 'zip' as the backup format", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide5, 
                        Inches(1), 
                        Inches(4.8), 
                        Inches(4), 
                        Inches(0.5), 
                        "5. Download the backup file", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_image(slide5, "images/export_database_diagram.png", Inches(5), Inches(3), width=Inches(4), height=Inches(2))



# Slide 6: Uploading the Database to the Odoo Upgrade Platform
slide6_layout = presentation.slide_layouts[1]
slide6 = presentation.slides.add_slide(slide6_layout)
slide6.shapes.title.text = "Uploading the Database to the Odoo Upgrade Platform"
add_rectangle_with_text(slide6, 
                        Inches(1), 
                        Inches(2),
                        Inches(4),
                        Inches(0.5), 
                        "1. Visit the Odoo Upgrade Platform (https://upgrade.odoo.com)\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))

add_rectangle_with_text(slide6, 
                        Inches(1), 
                        Inches(2.7), 
                        Inches(4), 
                        Inches(0.5), 
                        "2. Log in with your Odoo account\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide6, 
                        Inches(1), 
                        Inches(3.4), 
                        Inches(4), 
                        Inches(0.5), 
                        "3. Choose the target Odoo version\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide6, 
                        Inches(1), 
                        Inches(4.1), 
                        Inches(4), 
                        Inches(0.5), 
                        "4. Upload the previously exported database backup file\n", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))
add_rectangle_with_text(slide6, 
                        Inches(1), 
                        Inches(4.8), 
                        Inches(4), 
                        Inches(0.5), 
                        "5. Wait for the upgrade process to complete", 
                        fill_color=RGBColor(221, 235, 247), 
                        text_color=RGBColor(0, 0, 0))


add_image(slide6, "images/upload_database_diagram.png", Inches(5), Inches(3), width=Inches(4), height=Inches(2))

# Slide 7: Testing the Upgraded Database
slide7_layout = presentation.slide_layouts[1]
slide7 = presentation.slides.add_slide(slide7_layout)
slide7.shapes.title.text = "Testing the Upgraded Database"
slide7.placeholders[1].text = (
    "1. Download the upgraded database from the Odoo Upgrade Platform\n"
    "2. Restore the upgraded database in a test environment\n"
    "3. Test all critical business processes and custom modules\n"
    "4. Report any issues found during testing\n"
    "5. Repeat the process until all issues are resolved"
)
add_image(slide7, "images/testing_diagram.png", Inches(1), Inches(3), width=Inches(4), height=Inches(2))

# Slide 8: Migrating the Database to the Production Environment
slide8_layout = presentation.slide_layouts[1]
slide8 = presentation.slides.add_slide(slide8_layout)
slide8.shapes.title.text = "Migrating the Database to the Production Environment"
slide8.placeholders[1].text = (
    "1. Schedule downtime for the migration\n"
    "2. Backup the current production database\n"
    "3. Restore the upgraded and tested database in the production environment\n"
    "4. Update the Odoo configuration file\n"
    "5. Restart the Odoo server\n"
    "6. Verify functionality and perform sanity checks"
)
add_image(slide8, "images/migration_diagram.png", Inches(1), Inches(3), width=Inches(4), height=Inches(2))

# Save the presentation
current_datetime = datetime.now().strftime("%Y-%m-%d_%H%M%S")
presentation_filename = f"Odoo_Database_Upgrade_Process_{current_datetime}.pptx"

# Save the new presentation
presentation.save(presentation_filename)

